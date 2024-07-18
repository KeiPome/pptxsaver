import os
import sys
import datetime
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PPTX:
    def __init__(self,
                 filename,
                 title,
                 slide_title,
                 testname,
                 img_path,
                 slide_path
                 ):
        
        self.filename = filename
        self.title = title
        self.slide_title = slide_title
        self.img_path = img_path
        self.slide_path = slide_path
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.title_shape = self.slide.shapes.title
        self.title_shape.text = self.title
        self.title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        self.testname = testname

    def make_img_dir(self):
        if not os.path.exists(self.img_path):
            os.makedirs(self.img_path)
        else:
            shutil.rmtree(self.img_path)
            os.makedirs(self.img_path)
            
    def make_slide_dir(self):
        if not os.path.exists(self.slide_path):
            os.makedirs(self.slide_path)
        else:
            shutil.rmtree(self.slide_path)
            os.makedirs(self.slide_path)
    
    def save_img(self):
        filenames = []
        now = datetime.datetime.now()
        # yymmddhhmmss
        now_str = now.strftime('%Y%m%d')
        img_filename = self.img_path + '/' + self.testname + '_' + now_str + '.png'
        while os.path.exists(img_filename):
            n=1
            img_filename = f'{self.img_path}/{self.testname}_{now_str}-{str(n)}.png'
            n=+1
        plt.savefig(img_filename)
        filenames.append(img_filename)
        return img_filename
        
    def make_slide(self, files):
        for file in files:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            title_shape = slide.shapes.title
            title_shape.text = self.slide_title
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            slide.shapes.add_picture(file, Inches(1), Inches(1), Inches(8), Inches(6))
        self.prs.save(f'{self.slide_path}/{self.filename}.pptx')